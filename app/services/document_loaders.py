"""LangChain 文档加载器封装

提供统一的文档加载接口，支持多种文档格式。
使用 LangChain 原生加载器，简化封装层级。
"""

from dataclasses import dataclass
from enum import Enum
from pathlib import Path
from typing import Any

import httpx
import openpyxl
from bs4 import BeautifulSoup
from langchain_community.document_loaders import (
    PyMuPDFLoader,
    TextLoader,
    UnstructuredMarkdownLoader,
)
from langchain_core.documents import Document
from PIL import Image

from app.observability.logging import get_logger

logger = get_logger(__name__)


class DocumentFormat(str, Enum):
    """支持的文档格式"""

    PDF = "pdf"
    WORD = "word"
    EXCEL = "excel"
    PPT = "ppt"
    TEXT = "text"
    MARKDOWN = "markdown"
    WEB = "web"
    IMAGE = "image"

    @classmethod
    def from_extension(cls, ext: str) -> "DocumentFormat | None":
        """从文件扩展名获取格式

        Args:
            ext: 文件扩展名（包含点号，如 .pdf）

        Returns:
            DocumentFormat 实例或 None
        """
        ext = ext.lower()
        mapping = {
            ".pdf": cls.PDF,
            ".doc": cls.WORD,
            ".docx": cls.WORD,
            ".xls": cls.EXCEL,
            ".xlsx": cls.EXCEL,
            ".xlsm": cls.EXCEL,
            ".ppt": cls.PPT,
            ".pptx": cls.PPT,
            ".txt": cls.TEXT,
            ".md": cls.MARKDOWN,
            ".markdown": cls.MARKDOWN,
            # 图片格式
            ".png": cls.IMAGE,
            ".jpg": cls.IMAGE,
            ".jpeg": cls.IMAGE,
            ".gif": cls.IMAGE,
            ".bmp": cls.IMAGE,
            ".tiff": cls.IMAGE,
            ".webp": cls.IMAGE,
        }
        return mapping.get(ext)

    @classmethod
    def get_supported_extensions(cls) -> list[str]:
        """获取所有支持的文件扩展名"""
        return [
            ".pdf",
            ".doc",
            ".docx",
            ".xls",
            ".xlsx",
            ".xlsm",
            ".ppt",
            ".pptx",
            ".txt",
            ".md",
            ".markdown",
            ".png",
            ".jpg",
            ".jpeg",
            ".gif",
            ".bmp",
            ".tiff",
            ".webp",
        ]


@dataclass
class DocumentLoadResult:
    """文档加载结果

    兼容 LangChain Document，提供额外元数据。

    Attributes:
        content: 提取的文本内容
        metadata: 文档元数据
        page_count: 页数（如果适用）
        format: 文档格式
        source: 来源文件路径或 URL
    """

    content: str
    metadata: dict[str, Any]
    page_count: int
    format: DocumentFormat
    source: str

    def to_langchain_document(self) -> Document:
        """转换为 LangChain Document

        Returns:
            Document 实例
        """
        return Document(page_content=self.content, metadata=self.metadata)

    def to_dict(self) -> dict[str, Any]:
        """转换为字典"""
        return {
            "content": self.content,
            "metadata": self.metadata,
            "page_count": self.page_count,
            "format": self.format.value,
            "source": self.source,
        }


class PDFLoader:
    """PDF 文档加载器（使用 LangChain PyMuPDFLoader）

    封装 LangChain 的 PyMuPDFLoader，提供统一接口。
    """

    def __init__(self, source: str, extract_images: bool = False) -> None:
        """初始化 PDF 加载器

        Args:
            source: 文件路径
            extract_images: 是否提取图片（预留）
        """
        self.source = source
        self.format = DocumentFormat.PDF
        self.extract_images = extract_images

    async def load(self) -> DocumentLoadResult:
        """加载 PDF 文档

        Returns:
            DocumentLoadResult: 加载结果
        """
        try:
            # 使用 LangChain 的 PyMuPDFLoader
            loader = PyMuPDFLoader(
                self.source,
                extract_images=self.extract_images,
            )
            documents = loader.load()

            # 合并所有页面内容
            content_parts = []
            metadata = {"source": self.source, "pages": []}

            for i, doc in enumerate(documents):
                content_parts.append(doc.page_content)
                page_meta = {
                    "page": i + 1,
                    "page_content": doc.page_content[:200],  # 预览
                }
                # 复制其他元数据
                for key, value in doc.metadata.items():
                    if key not in page_meta:
                        page_meta[key] = value
                metadata["pages"].append(page_meta)

            content = "\n\n".join(content_parts)

            logger.info(
                "pdf_loaded",
                source=self.source,
                page_count=len(documents),
                content_length=len(content),
            )

            return DocumentLoadResult(
                content=content,
                metadata=metadata,
                page_count=len(documents),
                format=self.format,
                source=self.source,
            )

        except Exception as e:
            logger.error(
                "pdf_load_failed",
                source=self.source,
                error=str(e),
            )
            raise


class WordLoader:
    """Word 文档加载器（使用 python-docx）

    未来可替换为 LangChain 的 UnstructuredWordDocumentLoader。
    """

    async def load(self) -> DocumentLoadResult:
        """加载 Word 文档

        Returns:
            DocumentLoadResult: 加载结果
        """
        try:
            from docx import Document

            doc = Document(self.source)

            # 提取段落内容
            content_parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    content_parts.append(para.text)

            # 提取表格内容
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        content_parts.append(row_text)

            content = "\n\n".join(content_parts)

            metadata = {
                "source": self.source,
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables),
            }

            logger.info(
                "word_loaded",
                source=self.source,
                content_length=len(content),
            )

            return DocumentLoadResult(
                content=content,
                metadata=metadata,
                page_count=0,  # Word 文档不按页分
                format=self.format,
                source=self.source,
            )

        except Exception as e:
            logger.error(
                "word_load_failed",
                source=self.source,
                error=str(e),
            )
            raise


class ExcelLoader:
    """Excel 文档加载器（使用 openpyxl）

    未来可替换为 LangChain 的 UnstructuredExcelLoader。
    """

    async def load(self) -> DocumentLoadResult:
        """加载 Excel 文档

        Returns:
            DocumentLoadResult: 加载结果
        """
        try:
            wb = openpyxl.load_workbook(self.source, data_only=True)

            content_parts = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                content_parts.append(f"## 工作表: {sheet_name}")

                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(
                        str(cell) if cell is not None else "" for cell in row
                    )
                    if any(cell is not None for cell in row):
                        content_parts.append(row_text)

                content_parts.append("")  # 工作表之间空行

            content = "\n".join(content_parts)

            metadata = {
                "source": self.source,
                "sheet_count": len(wb.sheetnames),
                "sheets": wb.sheetnames,
            }

            logger.info(
                "excel_loaded",
                source=self.source,
                content_length=len(content),
            )

            return DocumentLoadResult(
                content=content,
                metadata=metadata,
                page_count=0,
                format=self.format,
                source=self.source,
            )

        except Exception as e:
            logger.error(
                "excel_load_failed",
                source=self.source,
                error=str(e),
            )
            raise


class PPTLoader:
    """PowerPoint 文档加载器（使用 python-pptx）

    未来可替换为 LangChain 的 UnstructuredPowerPointLoader。
    """

    async def load(self) -> DocumentLoadResult:
        """加载 PowerPoint 文档

        Returns:
            DocumentLoadResult: 加载结果
        """
        try:
            from pptx import Presentation

            prs = Presentation(self.source)

            content_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                content_parts.append(f"## 幻灯片 {slide_num}")

                # 提取文本框内容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_parts.append(shape.text)

                content_parts.append("")  # 幻灯片之间空行

            content = "\n".join(content_parts)

            metadata = {
                "source": self.source,
                "slide_count": len(prs.slides),
            }

            page_count = len(prs.slides)

            logger.info(
                "ppt_loaded",
                source=self.source,
                page_count=page_count,
                content_length=len(content),
            )

            return DocumentLoadResult(
                content=content,
                metadata=metadata,
                page_count=page_count,
                format=self.format,
                source=self.source,
            )

        except Exception as e:
            logger.error(
                "ppt_load_failed",
                source=self.source,
                error=str(e),
            )
            raise


class TextLoaderWrapper:
    """纯文本文档加载器（使用 LangChain TextLoader）

    封装 LangChain 的 TextLoader，提供统一接口。
    """

    async def load(self) -> DocumentLoadResult:
        """加载文本文档

        Returns:
            DocumentLoadResult: 加载结果
        """
        try:
            loader = TextLoader(self.source, encoding=self.encoding)
            documents = loader.load()

            content = documents[0].page_content if documents else ""
            metadata = {"source": self.source}
            if documents and documents[0].metadata:
                metadata.update(documents[0].metadata)

            logger.info(
                "text_loaded",
                source=self.source,
                content_length=len(content),
            )

            return DocumentLoadResult(
                content=content,
                metadata=metadata,
                page_count=0,
                format=self.format,
                source=self.source,
            )

        except Exception as e:
            logger.error(
                "text_load_failed",
                source=self.source,
                error=str(e),
            )
            raise


class MarkdownLoaderWrapper:
    """Markdown 文档加载器（使用 LangChain UnstructuredMarkdownLoader）

    封装 LangChain 的 UnstructuredMarkdownLoader，提供统一接口。
    """

    async def load(self) -> DocumentLoadResult:
        """加载 Markdown 文档

        Returns:
            DocumentLoadResult: 加载结果
        """
        try:
            loader = UnstructuredMarkdownLoader(self.source)
            documents = loader.load()

            content_parts = [doc.page_content for doc in documents]
            content = "\n\n".join(content_parts)

            metadata = {"source": self.source, "format": "markdown"}
            if documents and documents[0].metadata:
                metadata.update(documents[0].metadata)

            logger.info(
                "markdown_loaded",
                source=self.source,
                content_length=len(content),
            )

            return DocumentLoadResult(
                content=content,
                metadata=metadata,
                page_count=0,
                format=self.format,
                source=self.source,
            )

        except Exception as e:
            logger.error(
                "markdown_load_failed",
                source=self.source,
                error=str(e),
            )
            raise


class WebLoaderWrapper:
    """网页文档加载器（使用 httpx + BeautifulSoup）

    未来可替换为 LangChain 的 WebLoader。
    """

    async def load(self) -> DocumentLoadResult:
        """加载网页内容

        Returns:
            DocumentLoadResult: 加载结果
        """
        try:
            async with httpx.AsyncClient(timeout=self.timeout, follow_redirects=True) as client:
                response = await client.get(self.source)
                response.raise_for_status()

                # 解析 HTML
                soup = BeautifulSoup(response.text, "lxml")

                # 移除脚本和样式
                for script in soup(["script", "style", "nav", "footer", "header"]):
                    script.decompose()

                # 提取文本
                text = soup.get_text(separator="\n", strip=True)

                # 清理多余空行
                lines = [line.strip() for line in text.split("\n") if line.strip()]
                content = "\n\n".join(lines)

                # 提取元数据
                metadata = {
                    "source": self.source,
                    "format": "web",
                    "title": soup.title.string if soup.title else "",
                    "url": self.source,
                    "status_code": response.status_code,
                }

                # 提取 meta 描述
                description = soup.find("meta", attrs={"name": "description"})
                if description:
                    metadata["description"] = description.get("content", "")

                logger.info(
                    "web_loaded",
                    source=self.source,
                    content_length=len(content),
                )

                return DocumentLoadResult(
                    content=content,
                    metadata=metadata,
                    page_count=0,
                    format=self.format,
                    source=self.source,
                )

        except Exception as e:
            logger.error(
                "web_load_failed",
                source=self.source,
                error=str(e),
            )
            raise


class ImageLoaderWrapper:
    """图片加载器（基础实现）

    注意：当前版本仅提取图片基本信息，不执行 OCR。
    如需 OCR 功能，可集成 LangChain 的 UnstructuredImageLoader。
    """

    async def load(self) -> DocumentLoadResult:
        """加载图片信息

        Returns:
            DocumentLoadResult: 加载结果
        """
        try:
            with Image.open(self.source) as img:
                metadata = {
                    "source": self.source,
                    "format": img.format,
                    "mode": img.mode,
                    "size": img.size,
                    "width": img.width,
                    "height": img.height,
                }

                # 当前版本不执行 OCR，仅返回基本信息
                content = f"[Image: {img.format} {img.size[0]}x{img.size[1]}]"

                logger.info(
                    "image_loaded",
                    source=self.source,
                    size=img.size,
                )

                return DocumentLoadResult(
                    content=content,
                    metadata=metadata,
                    page_count=0,
                    format=self.format,
                    source=self.source,
                )

        except Exception as e:
            logger.error(
                "image_load_failed",
                source=self.source,
                error=str(e),
            )
            raise


class LoaderFactory:
    """文档加载器工厂

    根据文件类型选择合适的加载器。
    优先使用 LangChain 原生加载器，对于特殊格式提供适配层。
    """

    _loaders: dict[DocumentFormat, type] = {
        DocumentFormat.PDF: PDFLoader,
        DocumentFormat.WORD: WordLoader,
        DocumentFormat.EXCEL: ExcelLoader,
        DocumentFormat.PPT: PPTLoader,
        DocumentFormat.TEXT: TextLoaderWrapper,
        DocumentFormat.MARKDOWN: MarkdownLoaderWrapper,
        DocumentFormat.WEB: WebLoaderWrapper,
        DocumentFormat.IMAGE: ImageLoaderWrapper,
    }

    @classmethod
    def get_loader(
        cls,
        source: str,
        format: DocumentFormat | None = None,
        **kwargs: Any,
    ):
        """根据文件路径或格式获取加载器

        Args:
            source: 文件路径或 URL
            format: 文档格式（如果为 None，自动检测）
            **kwargs: 额外的加载器参数

        Returns:
            加载器实例（具体类型取决于文档格式）

        Raises:
            ValueError: 不支持的文件格式
        """
        if format is None:
            # 从文件扩展名检测格式
            path = Path(source)

            # URL 处理
            if source.startswith(("http://", "https://")):
                format = DocumentFormat.WEB
            else:
                ext = path.suffix.lower()
                format = DocumentFormat.from_extension(ext)

        if format is None:
            raise ValueError(
                f"Unsupported file format: {Path(source).suffix}. "
                f"Supported: {DocumentFormat.get_supported_extensions()}"
            )

        loader_class = cls._loaders.get(format)
        if loader_class is None:
            raise ValueError(f"No loader available for format: {format}")

        return loader_class(source, **kwargs)

    @classmethod
    def is_supported(cls, filename: str) -> bool:
        """检查文件是否支持

        Args:
            filename: 文件名

        Returns:
            是否支持
        """
        ext = Path(filename).suffix.lower()
        return ext in DocumentFormat.get_supported_extensions()

    @classmethod
    def get_supported_formats(cls) -> list[str]:
        """获取支持的格式列表

        Returns:
            格式名称列表
        """
        return [fmt.value for fmt in DocumentFormat]


async def load_document(
    source: str,
    format: DocumentFormat | None = None,
    **kwargs: Any,
) -> DocumentLoadResult:
    """便捷函数：加载文档

    Args:
        source: 文件路径或 URL
        format: 文档格式（自动检测）
        **kwargs: 额外的加载器参数

    Returns:
        DocumentLoadResult: 加载结果

    Raises:
        ValueError: 不支持的文件格式

    Examples:
        ```python
        # 从文件加载
        result = await load_document("path/to/document.pdf")

        # 从 URL 加载
        result = await load_document("https://example.com")

        # 指定格式
        result = await load_document("path/to/doc.pdf", format=DocumentFormat.PDF)
        ```
    """
    loader = LoaderFactory.get_loader(source, format, **kwargs)
    return await loader.load()


__all__ = [
    "DocumentFormat",
    "DocumentLoadResult",
    "PDFLoader",
    "WordLoader",
    "ExcelLoader",
    "PPTLoader",
    "TextLoaderWrapper",
    "MarkdownLoaderWrapper",
    "WebLoaderWrapper",
    "ImageLoaderWrapper",
    "LoaderFactory",
    "load_document",
]
